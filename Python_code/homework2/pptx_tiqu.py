#!/usr/bin/env python
# -*- coding: utf-8 -*-
# @Time    : 2024/5/20 20:06
# @Author  : 刘宇
# @File    : pptx_tiqu.py
import tkinter as tk
from tkinter import filedialog, messagebox
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os


class PPTExtractorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PPT内容提取器")

        # 设置窗口的初始位置和大小
        window_width = 600
        window_height = 110
        screen_width = root.winfo_screenwidth()
        screen_height = root.winfo_screenheight()
        x_coordinate = int((screen_width / 2) - (window_width / 2))
        y_coordinate = int((screen_height / 2) - (window_height / 2))
        self.root.geometry(f'{window_width}x{window_height}+{x_coordinate}+{y_coordinate}')

        # 加载自定义图标
        icon_path = 'pptLogo.png'  # 确保这里是正确的图标文件路径
        if os.path.exists(icon_path):
            icon = tk.PhotoImage(file=icon_path)
            # 如果需要调整图标大小，使用以下代码
            # 假设我们想要图标的大小为64x64像素
            original_width = icon.width()
            original_height = icon.height()
            new_width = 45
            new_height = 45
            if original_width > new_width or original_height > new_height:
                icon = icon.subsample(original_width // new_width, original_height // new_height)
            self.root.iconphoto(True, icon)  # 设置窗口图标
        else:
            print(f"图标文件未找到: {icon_path}")
        # 设置菜单栏
        menubar = tk.Menu(self.root)
        self.root.config(menu=menubar)

        filemenu = tk.Menu(menubar, tearoff=0)
        menubar.add_cascade(label="选择文件", menu=filemenu)
        filemenu.add_command(label="选择PPT文件", command=self.choose_ppt)
        filemenu.add_command(label="清除PPT路径", command=self.clear_ppt_path)
        # filemenu.add_separator()
        # filemenu.add_command(label="退出", command=self.root.quit)

        # 创建界面元素
        self.label_ppt_path = tk.Label(self.root, text="PPT文件路径: ")
        self.label_ppt_path.grid(row=0, column=0, sticky="e")

        self.entry_ppt_path = tk.Entry(self.root, width=50)
        self.entry_ppt_path.grid(row=0, column=1)

        # 创建一个框架用于容纳多选框，以便它们水平排布
        self.checkbox_frame = tk.Frame(self.root)
        self.checkbox_frame.grid(row=0, column=2, sticky="nsew")

        self.check_var_text = tk.BooleanVar()
        self.check_text = tk.Checkbutton(self.checkbox_frame, text="提取文字", variable=self.check_var_text)
        self.check_text.pack(side="left", anchor="w")

        self.check_var_images = tk.BooleanVar()
        self.check_images = tk.Checkbutton(self.checkbox_frame, text="提取图片", variable=self.check_var_images)
        self.check_images.pack(side="left", anchor="w")

        # 标签和输入框
        self.label_output_dir = tk.Label(self.root, text="输出目录路径: ")
        self.label_output_dir.grid(row=1, column=0, sticky="e")

        self.entry_output_dir = tk.Entry(self.root, width=50)
        self.entry_output_dir.grid(row=1, column=1)

        # 按钮放在输出目录路径输入框的右侧，并靠左对齐
        self.button_browse_output = tk.Button(self.root, text="选择输出目录", command=self.browse_output_dir)
        self.button_browse_output.grid(row=1, column=2, sticky="w")

        # 提取按钮
        self.button_extract = tk.Button(self.root, text="开始提取", command=self.extract_ppt,bg="#32993D")
        self.button_extract.grid(row=2, column=0, columnspan=3, pady=10)

    def choose_ppt(self):
        file_path = filedialog.askopenfilename(filetypes=[("PPT files", "*.pptx")])
        if file_path:
            self.entry_ppt_path.delete(0, tk.END)
            self.entry_ppt_path.insert(0, file_path)
        pass

    def clear_ppt_path(self):
        self.entry_ppt_path.delete(0, tk.END)
        pass

    def browse_output_dir(self):
        output_dir = filedialog.askdirectory()
        if output_dir:
            self.entry_output_dir.delete(0, tk.END)
            self.entry_output_dir.insert(0, output_dir)
        pass

        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def extract_ppt(self):
        ppt_path = self.entry_ppt_path.get()
        output_dir = self.entry_output_dir.get()

        if not ppt_path:
            messagebox.showerror("错误", "请先选择PPT文件")
            return

        if not output_dir:
            messagebox.showerror("错误", "请设置输出目录")
            return

        try:
            presentation = Presentation(ppt_path)
            if not os.path.exists(output_dir):
                os.makedirs(output_dir)
            self.extract_text_and_images(presentation, output_dir)

        except Exception as e:
            messagebox.showerror("错误", f"发生错误: {e}")
            return
        messagebox.showinfo("完成", "PPT提取完成！")
        self.root.destroy()

    def extract_text_and_images(self, presentation, output_dir):
        for slide_number, slide in enumerate(presentation.slides, start=1):
            # 提取文字
            if self.check_var_text.get():
                self.extract_text(slide, os.path.join(output_dir, f"word_{slide_number}.txt"))
            if self.check_var_images.get():
                self.extract_images(slide, output_dir, slide_number)

    def extract_text(self, slide, output_file_path):
        with open(output_file_path, 'w', encoding='utf-8') as file:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    file.write(shape.text)

    def extract_images(self, slide, output_dir, slide_number):
        image_counter = 1
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                if image:
                    image_bytes = image.blob
                    if image_bytes:
                        base_filename = f"image_{slide_number}"
                        image_filename = f"{base_filename}_{image_counter}.jpg"
                        image_filepath = os.path.join(output_dir, image_filename)
                        with open(image_filepath, 'wb') as img_file:
                            img_file.write(image_bytes)
                            print(f"图片提取成功: {image_filepath}")
                        image_counter += 1


if __name__ == "__main__":
    root = tk.Tk()
    app = PPTExtractorApp(root)
    root.mainloop()
