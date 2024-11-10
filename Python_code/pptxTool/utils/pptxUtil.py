import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.shapes.picture import Picture
from uuid import uuid1
import os


def select_file():
    root = tk.Tk()
    root.withdraw()  # 隐藏主窗口
    file_path = filedialog.askopenfilename(
        title="选择PPT文件",
        filetypes=[("PowerPoint Files", "*.pptx")]
    )
    return file_path


def extract_ppt_content(data_path, output_folder, extract_text=True, extract_image=True):
    ppt = Presentation(data_path)

    # 确保输出文件夹存在
    os.makedirs(output_folder, exist_ok=True)

    if extract_text:
        # 提取文本
        text_runs = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
        with open(os.path.join(output_folder, "extracted_text.txt"), "w", encoding="utf-8") as file:
            for text in text_runs:
                file.write(text + '\n')

    if extract_image:
        # 提取图片
        for slide in ppt.slides:
            for shape in slide.shapes:
                if isinstance(shape, Picture):
                    image_format = shape.image.content_type.split('/')[-1]
                    image_filename = f"{uuid1()}.{image_format}"
                    image_path = os.path.join(output_folder, image_filename)
                    with open(image_path, "wb") as f:
                        f.write(shape.image.blob)


def show_gui():
    root = tk.Tk()
    root.title("PPT内容提取工具")

    # 文本提取选项
    text_var = tk.BooleanVar()
    text_check = tk.Checkbutton(root, text="提取文本", variable=text_var)
    text_check.pack()

    # 图片提取选项
    image_var = tk.BooleanVar()
    image_check = tk.Checkbutton(root, text="提取图片", variable=image_var)
    image_check.pack()

    # 选择文件按钮和显示路径的文本框
    file_label = tk.Label(root, text="未选择文件")
    file_label.pack()

    select_button = tk.Button(root, text="选择PPT文件", command=lambda: process_file(text_var.get(), image_var.get(), file_label))
    select_button.pack()

    root.mainloop()


def process_file(extract_text, extract_image, file_label):
    data_path = select_file()
    if data_path:
        file_label.config(text=f"已选择文件: {data_path}")
        output_folder = '../output/'
        extract_ppt_content(data_path, output_folder, extract_text, extract_image)
        print(f"内容已成功提取并保存到 {output_folder}")
    else:
        file_label.config(text="未选择任何文件")
        print("未选择任何文件")


if __name__ == "__main__":
    show_gui()
